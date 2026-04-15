#!/usr/bin/env python3
"""
ingest.py — Load class materials into Pinecone
Supported: .pdf, .pptx, .docx, .txt, .md, .html, .csv

Usage:
  python ingest.py lecture1.pdf                        # single file
  python ingest.py lecture1.pdf notes.docx slides.pptx # multiple files
  python ingest.py materials.zip                       # zip of mixed files
  python ingest.py ./lectures/                         # folder of mixed files
"""

import os
import sys
import argparse
import tempfile
import zipfile
from pathlib import Path

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md", ".html", ".csv", ".json"}

EXTENSION_LABELS = {
    ".pdf":  "PDF",
    ".pptx": "PowerPoint",
    ".docx": "Word Doc",
    ".txt":  "Text",
    ".md":   "Markdown",
    ".html": "HTML",
    ".csv":  "CSV",
    ".json": "JSON",
}


def get_loader(file_path):
    """Return the appropriate LangChain loader for a given file."""
    from langchain_community.document_loaders import (
        PyMuPDFLoader,
        Docx2txtLoader,
        TextLoader,
        UnstructuredMarkdownLoader,
        BSHTMLLoader,
        CSVLoader,
    )
    from langchain_core.document_loaders import BaseLoader
    from langchain_core.documents import Document

    def _pdf_loader(file_path):
        """Return PyMuPDFLoader normally; fall back to OCR if the PDF is image-based."""
        import fitz
        doc = fitz.open(str(file_path))
        sample_text = "".join(doc[i].get_text() for i in range(min(3, doc.page_count)))
        doc.close()
        if len(sample_text.strip()) > 50:
            return PyMuPDFLoader(str(file_path))

        # Image-based scan — use Tesseract OCR
        class OcrPdfLoader(BaseLoader):
            def lazy_load(self):
                import pytesseract
                from pdf2image import convert_from_path
                images = convert_from_path(str(file_path), dpi=300)
                for i, img in enumerate(images, 1):
                    text = pytesseract.image_to_string(img)
                    if text.strip():
                        yield Document(page_content=text, metadata={"page": i})
        return OcrPdfLoader()

    class PptxLoader(BaseLoader):
        def __init__(self, path):
            self.path = path
        def lazy_load(self):
            from pptx import Presentation
            prs = Presentation(self.path)
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
                content = "\n".join(t for t in texts if t.strip())
                if content.strip():
                    yield Document(page_content=content, metadata={"slide": i})

    class CaseOpinionsLoader(BaseLoader):
        """Loads a case opinions JSON file (contracts-master.json schema)."""
        def __init__(self, path):
            self.path = path
        def lazy_load(self):
            import json, re
            with open(self.path, encoding="utf-8") as f:
                data = json.load(f)
            cases = data.get("cases", []) if isinstance(data, dict) else data
            for case in cases:
                name = case.get("case_name", "Unknown")
                citation = case.get("citation", "")
                court = case.get("court", "")
                date_filed = case.get("date_filed", "")
                header = f"{name} | {citation} | {court} | {date_filed}\n\n"
                for opinion in case.get("opinions", []):
                    raw = opinion.get("text", "")
                    # Strip XML/HTML tags
                    text = re.sub(r"<[^>]+>", " ", raw)
                    text = re.sub(r"\s{2,}", " ", text).strip()
                    if not text:
                        continue
                    yield Document(
                        page_content=header + text,
                        metadata={
                            "case_name": name or "",
                            "citation": citation or "",
                            "opinion_type": opinion.get("type") or "",
                            "date_filed": date_filed or "",
                        },
                    )

    ext = file_path.suffix.lower()
    loaders = {
        ".pdf":  lambda: _pdf_loader(file_path),
        ".pptx": lambda: PptxLoader(str(file_path)),
        ".docx": lambda: Docx2txtLoader(str(file_path)),
        ".txt":  lambda: TextLoader(str(file_path), encoding="utf-8"),
        ".md":   lambda: UnstructuredMarkdownLoader(str(file_path)),
        ".html": lambda: BSHTMLLoader(str(file_path)),
        ".csv":  lambda: CSVLoader(str(file_path)),
        ".json": lambda: CaseOpinionsLoader(str(file_path)),
    }
    return loaders[ext]()


def collect_files(inputs):
    """Resolve inputs (files, folders, zips) into a flat list of supported file paths."""
    file_paths = []
    temp_dirs = []

    for raw in inputs:
        path = Path(raw)

        if not path.exists():
            print(f"⚠️  Skipping (not found): {path}")
            continue

        # Supported file type
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            file_paths.append(path)

        # Zip — extract and find all supported files inside
        elif path.is_file() and path.suffix.lower() == ".zip":
            tmp = tempfile.mkdtemp()
            temp_dirs.append(tmp)
            print(f"📦 Extracting zip: {path.name}")
            with zipfile.ZipFile(path, "r") as zf:
                zf.extractall(tmp)
            found = [
                p for p in Path(tmp).rglob("*")
                if p.is_file()
                and p.suffix.lower() in SUPPORTED_EXTENSIONS
                and not p.name.startswith("._")
                and "__MACOSX" not in p.parts
            ]
            print(f"   Found {len(found)} supported file(s) inside")
            file_paths.extend(found)

        # Folder — find all supported files recursively
        elif path.is_dir():
            found = [
                p for p in path.rglob("*")
                if p.is_file()
                and p.suffix.lower() in SUPPORTED_EXTENSIONS
                and not p.name.startswith("._")
                and "__MACOSX" not in p.parts
            ]
            print(f"📁 Scanning folder '{path.name}': found {len(found)} supported file(s)")
            file_paths.extend(found)

        else:
            ext = path.suffix.lower()
            if path.is_file() and ext not in SUPPORTED_EXTENSIONS:
                print(f"⚠️  Skipping unsupported file type '{ext}': {path.name}")
                print(f"   Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")
            else:
                print(f"⚠️  Skipping: {path}")

    return file_paths, temp_dirs


def ingest_file(file_path, splitter, embeddings, index_name, pinecone_key):
    """Load, chunk, and upsert a single file. Returns chunk count."""
    from langchain_pinecone import PineconeVectorStore

    loader = get_loader(file_path)
    docs = loader.load()

    chunks = splitter.split_documents(docs)
    for chunk in chunks:
        chunk.metadata["source_file"] = file_path.name
        chunk.metadata["file_type"] = EXTENSION_LABELS.get(file_path.suffix.lower(), "Unknown")

    PineconeVectorStore.from_documents(
        chunks,
        embeddings,
        index_name=index_name,
        pinecone_api_key=pinecone_key,
    )
    return len(chunks)


def main():
    parser = argparse.ArgumentParser(description="Ingest class materials into Pinecone")
    parser.add_argument(
        "inputs",
        nargs="+",
        help="File(s), zip(s), or folder(s) — supports .pdf .pptx .docx .txt .md .html .csv",
    )
    parser.add_argument("--index", default="class-materials", help="Pinecone index name (default: class-materials)")
    parser.add_argument("--chunk-size", type=int, default=500, help="Chunk size (default: 500)")
    parser.add_argument("--chunk-overlap", type=int, default=50, help="Chunk overlap (default: 50)")
    args = parser.parse_args()

    # Check env
    pinecone_key = os.environ.get("PINECONE_API_KEY")
    if not pinecone_key:
        print("❌ Missing PINECONE_API_KEY. Set it with: export PINECONE_API_KEY=your-key")
        sys.exit(1)

    # Import dependencies
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        from langchain_huggingface import HuggingFaceEmbeddings
        from pinecone import Pinecone, ServerlessSpec
    except ImportError as e:
        print(f"❌ Missing dependency: {e}")
        print("   Run: pip install -r requirements.txt")
        sys.exit(1)

    # Collect all supported files from inputs
    file_paths, temp_dirs = collect_files(args.inputs)

    if not file_paths:
        print(f"❌ No supported files found.")
        print(f"   Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")
        sys.exit(1)

    # Show what we found, grouped by type
    from collections import Counter
    type_counts = Counter(EXTENSION_LABELS.get(p.suffix.lower(), p.suffix) for p in file_paths)
    print(f"\n📚 Found {len(file_paths)} file(s) to ingest:")
    for label, count in sorted(type_counts.items()):
        print(f"   • {count}x {label}")
    print()

    # Set up Pinecone index
    print(f"🌲 Connecting to Pinecone index '{args.index}'...")
    pc = Pinecone(api_key=pinecone_key)
    existing_indexes = [idx.name for idx in pc.list_indexes()]
    if args.index not in existing_indexes:
        print(f"   Creating new index '{args.index}'...")
        pc.create_index(
            name=args.index,
            dimension=384,
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1"),
        )
    print(f"   ✅ Ready\n")

    # Set up shared splitter + embeddings (load model once)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=args.chunk_size,
        chunk_overlap=args.chunk_overlap,
    )
    print("🔢 Loading embedding model (first run downloads ~80MB)...")
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    print("   ✅ Model ready\n")

    # Ingest each file
    total_chunks = 0
    failed = []

    for i, file_path in enumerate(file_paths, 1):
        label = EXTENSION_LABELS.get(file_path.suffix.lower(), file_path.suffix.upper())
        print(f"[{i}/{len(file_paths)}] {label}: {file_path.name}")
        try:
            n = ingest_file(file_path, splitter, embeddings, args.index, pinecone_key)
            total_chunks += n
            print(f"   ✅ {n} chunks uploaded")
        except Exception as e:
            print(f"   ❌ Failed: {e}")
            failed.append(file_path.name)

    # Cleanup temp dirs from zip extraction
    import shutil
    for tmp in temp_dirs:
        shutil.rmtree(tmp, ignore_errors=True)

    # Summary
    print(f"\n{'='*50}")
    print(f"🎉 Ingestion complete!")
    print(f"   ✅ {len(file_paths) - len(failed)} file(s) ingested")
    print(f"   📦 {total_chunks} total chunks in Pinecone")
    if failed:
        print(f"   ❌ {len(failed)} failed: {', '.join(failed)}")
    print(f"\n   Query with: python query.py \"your question here\"")


if __name__ == "__main__":
    main()
